#!/usr/bin/env python3
"""
Generate PowerPoint presentation from the Batch Inference System project.
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle

def add_content_slide(title, content_items):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    tf = content_shape.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = item
        else:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
        
        p.font.size = Pt(18)
        p.font.name = "Calibri"

def add_code_slide(title, code_text):
    """Add a slide with code example"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    
    # Add code
    code_top = Inches(1.5)
    code_height = Inches(5.5)
    code_box = slide.shapes.add_textbox(left, code_top, width, code_height)
    code_frame = code_box.text_frame
    code_frame.text = code_text
    code_frame.word_wrap = True
    
    for paragraph in code_frame.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.name = "Courier New"

# Slide 1: Title
add_title_slide(
    "Batch Inference System",
    "OpenAI-Style Batch Processing Platform\nProduction-ready • Kubernetes-native • Cost-optimized"
)

# Slide 2: Problem Statement
add_content_slide(
    "Problem Statement",
    [
        "Need to process large volumes of AI inference requests efficiently",
        "Balance cost optimization (spot instances) with reliability (dedicated)",
        "Ensure SLA compliance for time-sensitive workloads",
        "Provide visibility and control over batch processing",
        "",
        "Solution: Enterprise-grade batch inference platform"
    ]
)

# Slide 3: Key Features
add_content_slide(
    "Key Features",
    [
        "✓ JSONL File Upload & Batch Creation",
        "  • Upload large datasets via web UI or API",
        "  • Automatic splitting into individual requests",
        "",
        "✓ Priority-Based Scheduling",
        "  • Normal (1), Medium (5), High (10+) priority levels",
        "  • High-priority batches processed first",
        "",
        "✓ SLA-Aware Processing",
        "  • Automatic escalation from spot to dedicated workers",
        "  • Deadline-based GPU pool assignment",
        "",
        "✓ Request Deduplication",
        "  • Automatic detection of duplicate requests",
        "  • Reuse outputs from previous completions",
        "  • Significant cost and time savings"
    ]
)

# Slide 4: Architecture Overview
add_content_slide(
    "System Architecture",
    [
        "User/Portal → API Gateway → PostgreSQL",
        "                    ↓",
        "            Scheduler Service",
        "                    ↓",
        "        ┌───────────┴───────────┐",
        "        ↓                       ↓",
        "   Spot Workers          Dedicated Workers",
        "        ↓                       ↓",
        "    PostgreSQL ←───────────────┘",
        "        ↓",
        "  Batch Portal UI",
        "",
        "Key Components:",
        "• API Gateway (REST API)",
        "• Scheduler Service (batch orchestration)",
        "• GPU Workers (spot & dedicated pools)",
        "• PostgreSQL (durable queue + storage)",
        "• Batch Portal (web UI)",
        "• Monitoring Stack (Prometheus/Grafana)"
    ]
)

# Slide 5: Technology Stack
add_content_slide(
    "Technology Stack",
    [
        "Backend:",
        "• .NET 8 (C#)",
        "• ASP.NET Core (REST API)",
        "• Entity Framework Core (ORM)",
        "• PostgreSQL (database & queue)",
        "",
        "Infrastructure:",
        "• Kubernetes (orchestration)",
        "• Docker (containerization)",
        "• Prometheus (metrics)",
        "• Grafana (visualization)",
        "• Alertmanager (alerting)",
        "",
        "Frontend:",
        "• ASP.NET Razor Pages",
        "• Bootstrap 5",
        "• JavaScript (interactive UI)"
    ]
)

# Slide 6: Request Lifecycle
add_content_slide(
    "Request Lifecycle",
    [
        "1. Upload - User uploads JSONL file",
        "2. Create Batch - Batch created with priority level",
        "3. Deduplication - Check for duplicate inputs",
        "4. Queue - Requests queued in PostgreSQL",
        "5. Dequeue - Workers pull requests (priority-ordered)",
        "6. Process - GPU workers execute inference",
        "7. Complete - Results stored, batch finalized",
        "8. Download - User retrieves output file",
        "",
        "States: Queued → Running → Completed/Failed"
    ]
)

# Slide 7: Priority-Based Scheduling
add_content_slide(
    "Priority-Based Scheduling",
    [
        "How It Works:",
        "• Workers dequeue requests ordered by:",
        "  1. Batch Priority (highest first: 10+ → 5 → 1)",
        "  2. Creation Time (oldest first within same priority)",
        "",
        "Benefits:",
        "• Critical workloads processed first",
        "• Fair processing within priority levels",
        "• Independent of GPU pool selection",
        "• Cost optimization maintained",
        "",
        "Example:",
        "• Priority 10 batch (5 min old) → Processed first",
        "• Priority 1 batch (1 hour old) → Processed after high-priority"
    ]
)

# Slide 8: SLA-Aware Escalation
add_content_slide(
    "SLA-Aware Escalation",
    [
        "Default Strategy:",
        "• All requests start on spot workers (lower cost)",
        "• Monitor SLA deadline approaching",
        "• Automatically escalate to dedicated workers when needed",
        "",
        "Escalation Triggers:",
        "• Approaching completion deadline",
        "• Spot worker interruptions",
        "• Retry after failure",
        "",
        "Result:",
        "• Maximum cost savings (spot instances)",
        "• SLA compliance guaranteed",
        "• Automatic failover handling"
    ]
)

# Slide 9: Request Deduplication
add_content_slide(
    "Request Deduplication",
    [
        "How It Works:",
        "1. Compute SHA256 hash of normalized input payload",
        "2. Query database for completed requests with same hash",
        "3. If duplicate found:",
        "   • Copy output from original request",
        "   • Mark as deduplicated",
        "   • Complete immediately (no worker processing)",
        "",
        "Benefits:",
        "• Cost Savings - No redundant GPU processing",
        "• Time Savings - Instant completion for duplicates",
        "• Configurable Scope - Global or per-user deduplication",
        "",
        "Example:",
        "• 1000 requests, 200 duplicates → Only 800 processed",
        "• 20% cost reduction + faster completion"
    ]
)

# Slide 10: Database-Backed Queue
add_content_slide(
    "Database-Backed Queue",
    [
        "Why PostgreSQL?",
        "• No separate message broker needed",
        "• ACID transactions for reliability",
        "• FOR UPDATE SKIP LOCKED for safe concurrent dequeue",
        "• Persistent storage for all metadata",
        "",
        "Queue Semantics:",
        "• Atomic dequeue (transaction-based)",
        "• No double-processing",
        "• Priority ordering at database level",
        "• Automatic retry on failure",
        "",
        "Benefits:",
        "• Simpler architecture",
        "• Built-in persistence",
        "• Strong consistency guarantees"
    ]
)

# Slide 11: Batch Portal UI
add_content_slide(
    "Batch Portal UI",
    [
        "Dashboard:",
        "• Real-time statistics",
        "• Recent batches overview",
        "• System health indicators",
        "• Quick actions",
        "",
        "Batch Management:",
        "• Create batches with file upload",
        "• Priority selection",
        "• Filter, sort, and paginate",
        "• Cancel running batches",
        "• Clone existing batches",
        "",
        "Batch Details:",
        "• Visual progress tracking",
        "• Request-level details",
        "• Timeline visualization",
        "• Output download/preview",
        "• Error display with copy functionality"
    ]
)

# Slide 12: Monitoring & Observability
add_content_slide(
    "Monitoring & Observability",
    [
        "Prometheus Metrics:",
        "• Batch creation/completion rates",
        "• Request processing times",
        "• Worker pool utilization",
        "• SLA breach tracking",
        "• Deduplication statistics",
        "",
        "Grafana Dashboards:",
        "• Real-time system health",
        "• Performance trends",
        "• Cost optimization metrics",
        "• SLA compliance tracking",
        "",
        "Alertmanager:",
        "• SLA breach alerts",
        "• System health notifications",
        "• Worker pool capacity warnings"
    ]
)

# Slide 13: Key Differentiators
add_content_slide(
    "Key Differentiators",
    [
        "1. Production-Ready Architecture",
        "   • Kubernetes-native",
        "   • Scalable and resilient",
        "   • Full observability",
        "",
        "2. Cost Optimization",
        "   • Spot instance utilization",
        "   • Request deduplication",
        "   • Smart escalation only when needed",
        "",
        "3. Priority & SLA Management",
        "   • Multi-level priority system",
        "   • Automatic deadline tracking",
        "   • Guaranteed SLA compliance",
        "",
        "4. Developer Experience",
        "   • Clean API design",
        "   • Comprehensive documentation",
        "   • Easy local development"
    ]
)

# Slide 14: Use Cases
add_content_slide(
    "Use Cases",
    [
        "AI/ML Batch Processing:",
        "• Large-scale model inference",
        "• Batch predictions",
        "• Data transformation pipelines",
        "",
        "Cost-Sensitive Workloads:",
        "• Organizations needing spot instance savings",
        "• Variable processing requirements",
        "• Budget-conscious deployments",
        "",
        "Priority-Based Processing:",
        "• Mixed criticality workloads",
        "• Time-sensitive batch jobs",
        "• SLA-bound operations",
        "",
        "Deduplication Benefits:",
        "• Repeated processing of similar inputs",
        "• Caching frequently used queries",
        "• Cost reduction for common patterns"
    ]
)

# Slide 15: Performance Characteristics
add_content_slide(
    "Performance Characteristics",
    [
        "Scalability:",
        "• Horizontal scaling of workers",
        "• Database-backed queue handles high throughput",
        "• Kubernetes auto-scaling support",
        "",
        "Reliability:",
        "• ACID transactions",
        "• Automatic retry on failures",
        "• Spot interruption handling",
        "• Dead-letter queue support (future)",
        "",
        "Performance:",
        "• Priority-ordered processing",
        "• Efficient deduplication (indexed lookups)",
        "• Concurrent worker processing",
        "• Minimal overhead"
    ]
)

# Slide 16: Security & Multi-Tenancy
add_content_slide(
    "Security & Multi-Tenancy",
    [
        "Current:",
        "• User ID-based isolation",
        "• File-level access control",
        "• API authentication via headers",
        "",
        "Future Enhancements:",
        "• OAuth2/OIDC integration",
        "• JWT-based authentication",
        "• Row-level security per tenant",
        "• Audit logging",
        "• RBAC (Role-Based Access Control)"
    ]
)

# Slide 17: Deployment Options
add_content_slide(
    "Deployment Options",
    [
        "Local Development:",
        "• Docker Desktop with Kubernetes",
        "• Local PostgreSQL",
        "• Full feature parity",
        "",
        "Kubernetes:",
        "• Production-ready deployment",
        "• Helm charts available",
        "• Monitoring stack included",
        "• Persistent volumes",
        "",
        "Cloud Ready:",
        "• AWS EKS",
        "• Azure AKS",
        "• Google GKE",
        "• Any Kubernetes cluster"
    ]
)

# Slide 18: API Examples
code_examples = """Upload File:
POST /v1/files
Content-Type: multipart/form-data

Create Batch:
POST /v1/batches
{
  "inputFileId": "...",
  "metadata": {"priority": "10"}
}

Get Batch Status:
GET /v1/batches/{id}

Cancel Batch:
POST /v1/batches/{id}/cancel

Retry Request:
POST /v1/requests/{id}/retry"""
add_code_slide("API Examples", code_examples)

# Slide 19: Testing & Quality
add_content_slide(
    "Testing & Quality",
    [
        "Unit Tests:",
        "• View model mapping",
        "• API client behavior",
        "• Service logic",
        "• Schema validation",
        "",
        "Integration Tests:",
        "• End-to-end workflows",
        "• Database operations",
        "• Worker dequeue logic",
        "",
        "Test Infrastructure:",
        "• In-memory SQLite for fast tests",
        "• Mock services for isolation",
        "• Schema validation tests"
    ]
)

# Slide 20: Future Enhancements
add_content_slide(
    "Future Enhancements",
    [
        "Short Term:",
        "• Dead-letter queue implementation",
        "• Enhanced retry mechanisms",
        "• Cloud storage integration (S3/GCS)",
        "",
        "Medium Term:",
        "• Real GPU inference engine",
        "• Advanced authentication/authorization",
        "• Multi-region support",
        "",
        "Long Term:",
        "• Auto-scaling based on queue depth",
        "• Cost optimization analytics",
        "• Advanced scheduling algorithms"
    ]
)

# Slide 21: Getting Started
add_content_slide(
    "Getting Started",
    [
        "Prerequisites:",
        "• Docker Desktop with Kubernetes",
        "• .NET 8 SDK",
        "• kubectl",
        "",
        "Deploy:",
        "  ./scripts/redeploy-all.sh <tag>",
        "",
        "Access:",
        "• Portal: http://localhost:30081",
        "• API: http://localhost:30080",
        "• Grafana: http://localhost:30082",
        "",
        "Create Your First Batch:",
        "• Upload JSONL file",
        "• Select priority",
        "• Monitor progress"
    ]
)

# Slide 22: Metrics & KPIs
add_content_slide(
    "Metrics & KPIs",
    [
        "Throughput:",
        "• Batches processed per hour",
        "• Requests completed per minute",
        "• Deduplication rate",
        "",
        "Performance:",
        "• Average processing time",
        "• P95/P99 latencies",
        "• SLA compliance rate",
        "",
        "Cost:",
        "• Spot vs dedicated utilization",
        "• Cost savings from deduplication",
        "• Worker pool efficiency",
        "",
        "Reliability:",
        "• Success rate",
        "• Retry rate",
        "• Spot interruption frequency"
    ]
)

# Slide 23: Best Practices
add_content_slide(
    "Best Practices",
    [
        "Batch Creation:",
        "• Use appropriate priority levels",
        "• Group similar workloads",
        "• Monitor SLA deadlines",
        "",
        "Cost Optimization:",
        "• Leverage deduplication",
        "• Use spot instances when possible",
        "• Monitor escalation patterns",
        "",
        "Monitoring:",
        "• Set up Grafana dashboards",
        "• Configure alerts",
        "• Track key metrics",
        "",
        "Development:",
        "• Run tests before deployment",
        "• Use schema validation tests",
        "• Follow migration procedures"
    ]
)

# Slide 24: Conclusion
add_content_slide(
    "Conclusion",
    [
        "✓ Production-Ready batch inference platform",
        "✓ Cost-Optimized with spot instances and deduplication",
        "✓ SLA-Aware with automatic escalation",
        "✓ Priority-Based processing for critical workloads",
        "✓ Fully Observable with Prometheus/Grafana",
        "✓ Developer-Friendly with clean APIs and documentation",
        "",
        "Ready for:",
        "• Large-scale AI/ML batch processing",
        "• Cost-sensitive deployments",
        "• Priority-based workload management",
        "• Enterprise production use"
    ]
)

# Slide 25: Questions
add_title_slide(
    "Questions & Discussion",
    "Thank You!\n\nContact & Resources:\n• GitHub Repository\n• Documentation: README.md\n• API Documentation\n• Example Scripts"
)

# Save presentation
output_file = "Batch_Inference_System_Presentation.pptx"
prs.save(output_file)
print(f"✓ Presentation created: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")

